# -*- coding: utf-8 -*-
import sys, collections
from pptx import Presentation
from pptx.util import Emu, Pt

def hex_of(color):
    try:
        return str(color.rgb)
    except Exception:
        return None

def shape_fill_hex(shape):
    try:
        f = shape.fill
        if f.type is not None and f.type == 1:  # solid
            return hex_of(f.fore_color)
    except Exception:
        pass
    return None

def analyze(path):
    prs = Presentation(path)
    sw, sh = prs.slide_width, prs.slide_height
    print("="*80)
    print("FILE:", path)
    print("SLIDE SIZE: %.2f x %.2f in (EMU %d x %d)" % (sw/914400, sh/914400, sw, sh))
    print("NUM SLIDES:", len(prs.slides))
    fonts = collections.Counter()
    sizes = collections.Counter()
    fills = collections.Counter()
    yellowish = []
    bullets = 0
    for si, slide in enumerate(prs.slides, 1):
        for shp in slide.shapes:
            fh = shape_fill_hex(shp)
            if fh:
                fills[fh] += 1
                # detect yellow-ish band (high R,G low B)
                try:
                    r=int(fh[0:2],16); g=int(fh[2:4],16); b=int(fh[4:6],16)
                    if r>180 and g>150 and b<120:
                        txt = shp.text_frame.text[:30] if shp.has_text_frame else ""
                        yellowish.append((si, fh, round(shp.width/914400,2), round(shp.height/914400,2), round(shp.left/914400,2), round(shp.top/914400,2), txt))
                except Exception:
                    pass
            if shp.has_text_frame:
                for p in shp.text_frame.paragraphs:
                    if p.level and p.level>0:
                        bullets += 1
                    for r in p.runs:
                        nm = r.font.name
                        if nm: fonts[nm]+=1
                        if r.font.size: sizes[Pt(r.font.size.pt).pt]+=1
    print("\n-- TOP FONTS --")
    for k,v in fonts.most_common(10): print("  %-28s %d" % (k,v))
    print("-- TOP FONT SIZES (pt) --")
    for k,v in sorted(sizes.items(), key=lambda x:-x[1])[:12]: print("  %5.1fpt  x%d" % (k,v))
    print("-- TOP SHAPE FILL COLORS --")
    for k,v in fills.most_common(12): print("  #%-8s x%d" % (k,v))
    print("-- YELLOW-ISH BANDS (slide, color, w,h,left,top in, text) --")
    for y in yellowish[:15]: print("  ", y)

for p in sys.argv[1:]:
    try:
        analyze(p)
    except Exception as e:
        print("ERR", p, e)
