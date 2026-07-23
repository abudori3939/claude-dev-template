# -*- coding: utf-8 -*-
import sys
from pptx import Presentation
from pptx.util import Pt

def hx(c):
    try: return str(c.rgb)
    except: return None

def fillhex(s):
    try:
        if s.fill.type==1: return hx(s.fill.fore_color)
    except: pass
    return None

path=sys.argv[1]
slides=[int(x) for x in sys.argv[2:]] or [1,2]
prs=Presentation(path)
print("SIZE %.3f x %.3f in"%(prs.slide_width/914400, prs.slide_height/914400))
for sidx in slides:
    sl=prs.slides[sidx-1]
    print("\n#### SLIDE", sidx, "layout=", sl.slide_layout.name)
    for sh in sl.shapes:
        L=sh.left/914400 if sh.left is not None else None
        T=sh.top/914400 if sh.top is not None else None
        W=sh.width/914400 if sh.width is not None else None
        H=sh.height/914400 if sh.height is not None else None
        fontnm=fontsz=bold=algn=None
        txt=""
        if sh.has_text_frame:
            txt=sh.text_frame.text.replace("\n","/")[:40]
            for p in sh.text_frame.paragraphs:
                algn=p.alignment
                for r in p.runs:
                    fontnm=r.font.name;
                    if r.font.size: fontsz=r.font.size.pt
                    bold=r.font.bold
                    break
                if fontnm: break
        def f(x): return ("%.2f"%x) if isinstance(x,float) else str(x)
        print("  [%s] L%s T%s W%s H%s fill=%s font=%s sz=%s b=%s al=%s | %s"%(
            sh.shape_type, f(L),f(T),f(W),f(H), fillhex(sh), fontnm, fontsz, bold, algn, txt))
