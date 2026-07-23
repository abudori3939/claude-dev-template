# -*- coding: utf-8 -*-
"""既存資料の .pptx を複製し、中身を除去して空の templates/フォーマット.pptx を生成する。
   テーマ・スライドマスター・ロゴ・フォントを継承したまま、
   タイトルスライド + 空の本文スライドだけを残す。

   ⚠️ このスクリプトは「雛形」。KEEP 判定（どの要素を残すか）は体裁に依存するため、
   必ず自プロジェクトの agent_docs/project/format_rules.md の判別条件に合わせて調整すること。
"""
import copy, shutil, os
from pptx import Presentation

# ⚠️ 使用前に必ず書き換える ----------------------------------------------------
#   SRC        : 体裁の元にする既存資料 .pptx（プロジェクトごとに異なる）
#   KEEP_TEXT  : タイトルスライドで空にせず残す定数テキスト（例: 社名）。無ければ None
SRC = r"<体裁の元にする既存資料.pptx のフルパス>"
KEEP_TEXT = None  # 例: "株式会社○○"
# このスクリプトは agent_docs/stacks/pptx_deliverable/scripts/ にあるので、4つ上がリポジトリのルート
ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "..", ".."))
OUT = os.path.join(ROOT, r"templates", "フォーマット.pptx")
# ---------------------------------------------------------------------------

os.makedirs(os.path.dirname(OUT), exist_ok=True)
shutil.copyfile(SRC, OUT)

prs = Presentation(OUT)


def fill_hex(s):
    try:
        if s.fill.type == 1:
            return str(s.fill.fore_color.rgb)
    except Exception:
        pass
    return None


def blank_text(shape):
    """text_frameのrunテキストを空にし、書式(フォント/サイズ/太字)は維持する。"""
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.text = ""


def remove_shape(shape):
    shape._element.getparent().remove(shape._element)


# --- スライド削除: 1(タイトル) と 2(本文) のみ残す ---
# 元資料の「タイトル」「代表的な本文」が 1・2 枚目でない場合は、先に該当スライドを先頭に
# 移動するか、下の ids のインデックスを調整すること。
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
for sid in ids[2:]:
    sldIdLst.remove(sid)

slides = list(prs.slides)
title_slide, body_slide = slides[0], slides[1]

# --- スライド1(タイトル): 可変テキストを空に。KEEP_TEXT(定数)は残す ---
for sh in title_slide.shapes:
    if not sh.has_text_frame:
        continue
    if KEEP_TEXT and KEEP_TEXT in sh.text_frame.text:  # 定数（社名など）は残す
        continue
    blank_text(sh)  # タイトル・日付などを空に

# --- スライド2(本文): 共通レイアウト要素のみ残し、中身の図解・本文固有の図形は削除 ---
# ↓↓ 以下の KEEP 判定は一例（format_rules.md の判別条件に合わせて書き換える）↓↓
#   例: ロゴ(PICTURE) / 全幅の区切り線 / ページ番号(PLACEHOLDER) /
#       固有要素(特定の塗り色・全幅) / タイトル(28pt) / 本文(24pt・全幅クラス)
to_remove = []
for sh in body_slide.shapes:
    st = str(sh.shape_type)
    keep = False
    if "PICTURE" in st:  # ロゴ
        keep = True
    elif "LINE" in st and sh.width and sh.width / 914400 > 10:  # 区切り線(全幅)
        keep = True
    elif "PLACEHOLDER" in st:  # ページ番号
        keep = True
    elif fill_hex(sh) == "FFFF00" and sh.width and sh.width / 914400 > 10:
        # 固有要素の例(全幅の帯)。塗り色は format_rules.md に合わせる
        keep = True
        blank_text(sh)
    elif sh.has_text_frame:
        f = None
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size:
                    f = r.font.size.pt
                break
            if f:
                break
        w = sh.width / 914400 if sh.width else 0
        if f == 28.0:  # タイトル(サイズは format_rules.md に合わせる)
            keep = True
            blank_text(sh)
        elif f == 24.0 and w > 10:  # 本文(全幅クラス)
            keep = True
            blank_text(sh)
    if not keep:
        to_remove.append(sh)

for sh in to_remove:
    remove_shape(sh)

prs.save(OUT)
print("SAVED:", OUT)

# 確認: 生成物を再読込し、残った要素を一覧表示（このあと必ず目視確認もする）
chk = Presentation(OUT)
print("slides:", len(chk.slides))
for i, sl in enumerate(chk.slides, 1):
    print("-- slide", i, "--")
    for sh in sl.shapes:
        L = sh.left / 914400 if sh.left is not None else 0
        T = sh.top / 914400 if sh.top is not None else 0
        print("   %-16s L%.2f T%.2f fill=%s text=%r" % (
            str(sh.shape_type)[:16], L, T, fill_hex(sh),
            sh.text_frame.text[:20] if sh.has_text_frame else ""))
